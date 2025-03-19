import pandas as pd
from pptx import Presentation
import pdfplumber
import json

class DataProcessor:
    def __init__(self, json=None, csv=None, pdf=None, pptx=None):
        self.json = json.get_data() if json else []
        self.csv = csv.get_data() if csv else []
        self.pdf = pdf.get_data() if pdf else []
        self.pptx = pptx.get_data() if pptx else []
        self.data = []

        if json:
            self.data.append(self.json)
        if csv:
            self.data.append(self.csv)
        if pdf:
            self.data.append(self.pdf)
        if pptx:
            self.data.append(self.pptx)

    def get_data(self):
        return self.data
    
    def get_data_file_type(self, filetype):
        if filetype == 'json':
            return self.json
        if filetype == 'csv':
            return self.csv
        if filetype == 'pdf':
            return self.pdf
        if filetype == 'pptx':
            return self.pptx
        
    def set_file(self, json=None, csv=None, pdf=None, pptx=None):
        self.json = json.get_data() if json else []
        self.csv = csv.get_data() if csv else []
        self.pdf = pdf.get_data() if pdf else []
        self.pptx = pptx.get_data() if pptx else []
        self.data = []

        if json:
            self.data.append(self.json)
        if csv:
            self.data.append(self.csv)
        if pdf:
            self.data.append(self.pdf)
        if pptx:
            self.data.append(self.pptx)




class JSONData:
    def __init__(self, path=None):
        self.path = path
    
    def get_data(self):
        df = pd.read_json(self.path)
        self.df = df.to_dict(orient='records')
        return self.df

    
    def set_data(self, path=None):
        self.path = path

    

class CSVData:
    def __init__(self, path=None):
        self.path = path
    
    def get_data(self):
        df = pd.read_csv(self.path)
        self.df = df
        self.json_data = df.to_dict(orient="records")
        return self.json_data

    
    def set_data(self, path=None):
        self.path = path


class PDFData:
    def __init__(self, path=None):
        self.path = path
    
    def get_data(self):
        with pdfplumber.open(self.path) as pdf:
            pdf_text = ""
            for page in pdf.pages:
                pdf_text += page.extract_text()


        lines = pdf_text.strip().split("\n")
        self.json_data = []
        for line in lines[2:]:
            values = line.split(" ")

            self.json_data.append({
                "Year": values[0],
                "Quarter": values[1],
                "Revenue (in $)": values[2],
                "Memberships Sold": int(values[3]),
                "Ave Duration (Minutes)": int(values[4])
            })
        return self.json_data

    
    def set_data(self, path=None):
        self.path = path


    

class PPTXData:
    def __init__(self, path=None):
        self.path = path
    
    def get_data(self):
        presentation = Presentation(self.path)
        slides_data = []
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = ""
            header = ""
            slide_text_object = {}
            sub_header = ""

            for i, shape in enumerate(slide.shapes, start=0):
                if hasattr(shape, "text"):
                    if i == 0:
                        header = shape.text.strip()
                    else:
                        lines = shape.text.split("\n")
                        for j, line in enumerate(lines, start=0):
                            if ":" in line:
                                key, value = map(str.strip, line.split(":", 1))
                                if j == 0 and sub_header == "":
                                    sub_header = key
                                slide_text_object[key] = value
                    

                if shape.has_table:
                    table = shape.table
                    row_data = [""]
                    clmns = 0
                    clm_headers = [""]
                    table_data = []
                    sub_header = "table"

                    for k, row in enumerate(table.rows, start=0):
                        if k == 0:
                            clm_headers = [cell.text.strip() for cell in row.cells]
                            continue

                        row_data = [cell.text.strip() for cell in row.cells]
                        row_dict = {clm_headers[l]: row_data[l] for l in range(len(clm_headers))}
                        table_data.append(row_dict)
                    slide_text_object["Table Data"] = table_data
                    

            slides_data.append({
                "title": header,
                sub_header: slide_text_object
            })
        self.slides_data = slides_data

        return self.slides_data

    
    def set_data(self, path=None):
        self.path = path

    

if __name__ == "__main__":
    json_data = JSONData(path="../data/dataset1.json")
    json_data.get_data()
    csv_data = CSVData(path="../data/dataset2.csv")
    csv_data.get_data()
    pdf_data = PDFData(path="../data/dataset3.pdf")
    pdf_data.get_data()
    pptx_data = PPTXData(path="../data/dataset4.pptx")
    pptx_data.get_data()
    data_processor = DataProcessor(json_data, csv_data, pdf_data, pptx_data)
    data_processor.get_data()